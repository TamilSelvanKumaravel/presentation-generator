"""
Gamma.ai CORRECT Implementation - Light Backgrounds with Colored Boxes
Based on actual Gamma.ai slides analysis.
"""
import os
import random
from pathlib import Path
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from src.models.schemas import LLMResponse
from src.utils.logger import app_logger
from src.services.media.image_service import ImageService

# Gamma.ai Color Palette (from real slides)
GAMMA_COLORS = {
    "primary_blue": (10, 92, 140),      # #0A5C8C - Dark blue for text
    "light_blue_box": (212, 233, 247),  # #D4E9F7 - Content boxes
    "background": (245, 247, 250),      # #F5F7FA - Light gray background
    "text_gray": (107, 114, 128),       # #6B7280 - Body text
    "white": (255, 255, 255),           # #FFFFFF - Cards
    "accent_blue": (59, 130, 246),      # #3B82F6 - Accents
}

class PPTXGenerator:
    """Generate Gamma.ai-style presentations with light backgrounds and colored boxes."""
    
    def __init__(self, output_dir: str = "presentations"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(exist_ok=True)
        self.image_service = ImageService()
        app_logger.info("Gamma.ai Generator initialized (Light theme)")

    def _get_optimized_font_size(self, text: str, max_len: int, base_size: int, min_size: int = 14) -> Pt:
        """Calculate optimized font size based on text length."""
        length = len(text)
        if length <= max_len:
            return Pt(base_size)
        elif length <= max_len * 1.5:
            return Pt(int(base_size * 0.8))
        elif length <= max_len * 2:
            return Pt(int(base_size * 0.65))
        else:
            return Pt(max(min_size, int(base_size * 0.5)))
    
    def generate(self, llm_response: LLMResponse, style: str = "professional") -> str:
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Title slide
            title_bg = f"{llm_response.title} professional background"
            self._create_title_slide(prs, llm_response.title, llm_response.summary, title_bg)
            
            # Content slides - 3 layout types
            for idx, slide_data in enumerate(llm_response.slides):
                has_image = bool(getattr(slide_data, 'image_query', None))
                
                if has_image:
                    if idx % 3 == 0:
                        self._create_content_with_boxes(prs, slide_data)
                    elif idx % 3 == 1:
                        self._create_split_with_image(prs, slide_data)
                    else:
                        self._create_grid_layout(prs, slide_data)
                else:
                    self._create_simple_content(prs, slide_data)
            
            filename = self._generate_filename(llm_response.title)
            filepath = self.output_dir / filename
            prs.save(str(filepath))
            
            app_logger.info(f"Presentation saved: {filepath}")
            return str(filepath)
            
        except Exception as e:
            app_logger.error(f"Error generating PPTX: {str(e)}")
            raise
    
    def _add_light_background(self, slide, image_query=None):
        """Add light background with optional subtle image."""
        if image_query:
            image_path = self.image_service.get_image(image_query)
            if image_path:
                try:
                    # Add image
                    slide.shapes.add_picture(
                        image_path, Inches(0), Inches(0),
                        width=Inches(10), height=Inches(7.5)
                    )
                    # Add white overlay for readability
                    overlay = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(0), Inches(0), Inches(10), Inches(7.5)
                    )
                    overlay.fill.solid()
                    overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
                    overlay.fill.fore_color.brightness = 0.85  # 85% white
                    overlay.line.fill.background()
                    return
                except:
                    pass
        
        # Fallback: solid light background
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*GAMMA_COLORS["background"])
        bg.line.fill.background()
    
    def _create_title_slide(self, prs, title, summary, bg_query):
        """Title slide with light background and dark text - Gamma.ai style."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add solid light background first
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(7.5)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*GAMMA_COLORS["background"])
        bg.line.fill.background()
        
        # Add image as SUBTLE background on right side with high transparency
        image_path = self.image_service.get_image(bg_query)
        if image_path:
            try:
                # Add image
                pic = slide.shapes.add_picture(
                    image_path,
                    Inches(5), Inches(0.5),
                    width=Inches(4.5), height=Inches(6.5)
                )
                # Overlay removed to keep image visible
            except:
                pass
        
        # Title - large dark blue text (left side only)
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.0), Inches(4.0), Inches(2.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Segoe UI"
        p.font.name = "Segoe UI"
        # Dynamic font sizing
        p.font.size = self._get_optimized_font_size(title, 30, 44, 24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*GAMMA_COLORS["primary_blue"])
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.1
        
        # Subtitle - positioned BELOW title with proper spacing
        sub_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(3.8), Inches(4.0), Inches(2.8)
        )
        tf = sub_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.text = summary
        p.font.name = "Segoe UI"
        p.font.size = Pt(16)  # Reduced from 18pt
        p.font.color.rgb = RGBColor(*GAMMA_COLORS["primary_blue"])
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.4
        
        # Author/subtitle - at bottom
        author_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(6.8), Inches(4.5), Inches(0.5)
        )
        tf = author_box.text_frame
        p = tf.paragraphs[0]
        p.text = "by AI Research Team"
        p.font.name = "Segoe UI"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(*GAMMA_COLORS["text_gray"])
        p.alignment = PP_ALIGN.LEFT
    
    def _create_content_with_boxes(self, prs, slide_data):
        """Content slide with light blue boxes (like Gamma.ai Slide 3)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Light background
        self._add_light_background(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.title
        p.font.name = "Segoe UI"
        p.font.name = "Segoe UI"
        p.font.size = self._get_optimized_font_size(slide_data.title, 40, 30, 18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*GAMMA_COLORS["primary_blue"])
        p.alignment = PP_ALIGN.LEFT
        
        # Content boxes (left side)
        box_y = 1.8
        
        # Calculate uniform font size for all boxes
        content_items = slide_data.content[:2]
        font_sizes = [self._get_optimized_font_size(item, 50, 16, 10).pt for item in content_items]
        uniform_size = Pt(min(font_sizes)) if font_sizes else Pt(16)
        
        for i, point in enumerate(content_items):  # Max 2 boxes
            # Light blue box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), Inches(box_y), Inches(4.5), Inches(1.8)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*GAMMA_COLORS["light_blue_box"])
            box.line.fill.background()
            box.adjustments[0] = 0.1
            
            # Text in box
            text_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(box_y + 0.3), Inches(3.8), Inches(1.2)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = point
            p.font.name = "Segoe UI"
            p.font.size = uniform_size
            p.font.color.rgb = RGBColor(*GAMMA_COLORS["primary_blue"])
            p.line_spacing = 1.3
            
            box_y += 2.2
        
        # Images on right side
        image_path = self.image_service.get_image(slide_data.image_query)
        if image_path:
            try:
                # Top image
                slide.shapes.add_picture(
                    image_path,
                    Inches(5.8), Inches(1.5),
                    width=Inches(3.5), height=Inches(2.2)
                )
                # Bottom image (same for now)
                slide.shapes.add_picture(
                    image_path,
                    Inches(5.8), Inches(4),
                    width=Inches(3.5), height=Inches(2.2)
                )
            except:
                pass
    
    def _create_split_with_image(self, prs, slide_data):
        """Split layout with image and text boxes."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Light background
        self._add_light_background(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.title
        p.font.name = "Segoe UI"
        p.font.name = "Segoe UI"
        p.font.size = self._get_optimized_font_size(slide_data.title, 35, 36, 20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*GAMMA_COLORS["primary_blue"])
        p.alignment = PP_ALIGN.LEFT
        
        # Left side - content boxes
        box_y = 2
        
        # Calculate uniform font size
        content_items = slide_data.content[:3]
        font_sizes = [self._get_optimized_font_size(item, 60, 14, 10).pt for item in content_items]
        uniform_size = Pt(min(font_sizes)) if font_sizes else Pt(14)
        
        for i, point in enumerate(content_items):
            # Light blue box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.8), Inches(box_y), Inches(4.2), Inches(1.3)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*GAMMA_COLORS["light_blue_box"])
            box.line.fill.background()
            box.adjustments[0] = 0.1
            
            # Text
            text_box = slide.shapes.add_textbox(
                Inches(1.2), Inches(box_y + 0.25), Inches(3.5), Inches(0.8)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = point
            p.font.name = "Segoe UI"
            p.font.size = uniform_size
            p.font.color.rgb = RGBColor(*GAMMA_COLORS["primary_blue"])
            
            box_y += 1.5
        
        # Right side - image
        image_path = self.image_service.get_image(slide_data.image_query)
        if image_path:
            try:
                slide.shapes.add_picture(
                    image_path,
                    Inches(5.5), Inches(2),
                    width=Inches(3.8), height=Inches(4.5)
                )
            except:
                pass
    
    def _create_grid_layout(self, prs, slide_data):
        """Grid layout with 4 boxes (SWOT-style)."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Light background
        self._add_light_background(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.8)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.title
        p.font.name = "Segoe UI"
        p.font.name = "Segoe UI"
        p.font.size = self._get_optimized_font_size(slide_data.title, 35, 36, 20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*GAMMA_COLORS["primary_blue"])
        p.alignment = PP_ALIGN.CENTER
        
        # 2x2 grid of boxes
        positions = [
            (0.8, 2, "S"),
            (5.1, 2, "W"),
            (0.8, 4.5, "T"),
            (5.1, 4.5, "O")
        ]
        
        # Calculate uniform font size
        content_items = slide_data.content[:4]
        font_sizes = [self._get_optimized_font_size(item, 80, 12, 9).pt for item in content_items]
        uniform_size = Pt(min(font_sizes)) if font_sizes else Pt(12)
        
        for i, (x, y, letter) in enumerate(positions):
            if i >= len(slide_data.content):
                break
            
            # Light blue box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(4), Inches(2.2)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*GAMMA_COLORS["light_blue_box"])
            box.line.fill.background()
            box.adjustments[0] = 0.1
            
            # Letter in top-left corner
            letter_box = slide.shapes.add_textbox(
                Inches(x + 0.2), Inches(y + 0.1), Inches(0.6), Inches(0.6)
            )
            tf = letter_box.text_frame
            p = tf.paragraphs[0]
            p.text = letter
            p.font.name = "Segoe UI"
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(*GAMMA_COLORS["accent_blue"])
            p.alignment = PP_ALIGN.LEFT
            
            # Text - moved down to avoid overlap
            text_box = slide.shapes.add_textbox(
                Inches(x + 0.3), Inches(y + 0.6), Inches(3.4), Inches(1.4)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data.content[i]
            p.font.name = "Segoe UI"
            p.font.size = uniform_size
            p.font.color.rgb = RGBColor(*GAMMA_COLORS["primary_blue"])
            p.line_spacing = 1.2
    
    def _create_simple_content(self, prs, slide_data):
        """Simple content slide with boxes."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Light background
        self._add_light_background(slide)
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.8), Inches(8.4), Inches(1)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.title
        p.font.name = "Segoe UI"
        p.font.name = "Segoe UI"
        p.font.size = self._get_optimized_font_size(slide_data.title, 40, 28, 18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*GAMMA_COLORS["primary_blue"])
        p.alignment = PP_ALIGN.LEFT
        
        # Content boxes
        box_y = 2.5
        
        # Calculate uniform font size
        font_sizes = [self._get_optimized_font_size(item, 70, 16, 11).pt for item in slide_data.content]
        uniform_size = Pt(min(font_sizes)) if font_sizes else Pt(16)
        
        for point in slide_data.content:
            # Light blue box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(1.5), Inches(box_y), Inches(7), Inches(1)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*GAMMA_COLORS["light_blue_box"])
            box.line.fill.background()
            box.adjustments[0] = 0.1
            
            # Text
            text_box = slide.shapes.add_textbox(
                Inches(2), Inches(box_y + 0.2), Inches(6), Inches(0.6)
            )
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {point}"
            p.font.name = "Segoe UI"
            p.font.size = uniform_size
            p.font.color.rgb = RGBColor(*GAMMA_COLORS["primary_blue"])
            
            box_y += 1.2
    
    def _generate_filename(self, title):
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        safe_title = "".join(c if c.isalnum() or c in (' ', '-', '_') else '' for c in title)
        safe_title = safe_title.replace(' ', '_')[:20]
        return f"{safe_title}_gamma_{timestamp}.pptx"
