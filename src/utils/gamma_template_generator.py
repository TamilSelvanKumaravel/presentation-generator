"""
Simplified Gamma.ai-inspired template generator.
Creates clean layouts with proper styling - decorative elements added during generation.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

def create_gamma_template():
    """Create a clean Gamma-inspired template."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Set dark background on master
    master = prs.slide_master
    background = master.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42)  # Slate 900
    
    # Configure Title Layout (Index 0)
    title_layout = prs.slide_layouts[0]
    title_placeholder = title_layout.placeholders[0]
    title_placeholder.top = Inches(2.5)
    title_placeholder.left = Inches(1.5)
    title_placeholder.width = Inches(7)
    title_placeholder.height = Inches(1.5)
    
    # Configure subtitle if exists
    if len(title_layout.placeholders) > 1:
        subtitle_placeholder = title_layout.placeholders[1]
        subtitle_placeholder.top = Inches(4.2)
        subtitle_placeholder.left = Inches(1.5)
        subtitle_placeholder.width = Inches(7)
        subtitle_placeholder.height = Inches(1)
    
    # Configure Content Layout (Index 1)
    content_layout = prs.slide_layouts[1]
    content_title = content_layout.placeholders[0]
    content_title.top = Inches(0.5)
    content_title.left = Inches(0.5)
    content_title.width = Inches(9)
    content_title.height = Inches(0.7)
    
    # Configure body
    if len(content_layout.placeholders) > 1:
        content_body = content_layout.placeholders[1]
        content_body.top = Inches(1.8)
        content_body.left = Inches(1)
        content_body.width = Inches(8)
        content_body.height = Inches(4.7)
    
    # Save template
    output_path = Path("templates/theme_gamma.pptx")
    output_path.parent.mkdir(exist_ok=True)
    prs.save(str(output_path))
    print(f"✓ Created Gamma template: {output_path}")
    return output_path

if __name__ == "__main__":
    create_gamma_template()
    print("✓ Template generation complete!")
